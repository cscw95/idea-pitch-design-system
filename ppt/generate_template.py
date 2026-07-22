# -*- coding: utf-8 -*-
"""
Idea Pitch PPT Template — Design System Generator
DESIGN_SYSTEM.md 의 토큰/컴포넌트 정의를 python-pptx 로 구현한다.
실행: python3 generate_template.py  →  idea_pitch_template.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.oxml.ns import qn

# ──────────────────────────────── 1. 디자인 토큰 ────────────────────────────────
NAVY    = RGBColor(0x1F, 0x38, 0x64)
BLUE    = RGBColor(0x12, 0x73, 0xC4)
SKY     = RGBColor(0xDC, 0xE7, 0xF2)
INK     = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_TX = RGBColor(0x7F, 0x7F, 0x7F)
STEEL   = RGBColor(0x8C, 0x96, 0xA0)
BORDER  = RGBColor(0xC9, 0xCD, 0xD3)
RED     = RGBColor(0xC0, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# 기본 서체: 나눔스퀘어 (https://hangeul.naver.com/fonts/search?f=nanum)
# 웨이트별 패밀리가 분리되어 있으므로 bold 플래그 대신 weight 토큰으로 매핑한다.
# _ac 계열 = 한글 11,172자 전체 지원판
FONT_MAP = {
    'l':  "NanumSquareOTF_ac Light",
    'r':  "NanumSquareOTF_ac",
    'b':  "NanumSquareOTF_ac Bold",
    'eb': "NanumSquareOTF_ac ExtraBold",
}

MARGIN_X   = 0.45
CONTENT_W  = 12.43
ZONE_HEADLINE_Y = 1.10
ZONE_BODY_Y     = 1.95
ZONE_BODY_H     = 4.65
ZONE_TAKEAWAY_Y = 6.78

# ──────────────────────────────── 2. 저수준 헬퍼 ────────────────────────────────

def _ea(run, name):
    """한글(East Asian) 서체 지정 — font.name 은 라틴 서체만 설정하므로 별도 처리"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def _tf(shape, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    return tf


def para(tf, runs, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0, line=None, first=False, weight=None):
    """runs: str 또는 [(text, {size,bold,color,weight}), ...]
    weight: 'l'/'r'/'b'/'eb' — 나눔스퀘어 웨이트 토큰. 미지정 시 bold 여부로 r/b 자동 선택"""
    if first and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, opt in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        w = opt.get('weight', weight) or ('b' if opt.get('bold', bold) else 'r')
        fname = FONT_MAP[w]
        f.name = fname
        _ea(r, fname)
        f.size = Pt(opt.get('size', size))
        f.bold = False   # 웨이트는 패밀리로 표현 — 가짜 볼드 중복 방지
        f.color.rgb = opt.get('color', color)
    return p


def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=None, dash=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w)
        if dash:
            sp.line.dash_style = dash
    return sp


def hline(slide, x1, y, x2, color=INK, width=1.5, dash=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y), Inches(x2), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    ln.shadow.inherit = False
    if dash:
        ln.line.dash_style = dash
    return ln

# ──────────────────────────────── 3. 컴포넌트 ────────────────────────────────

def add_header(slide, deck_name, section_label):
    tb = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(0.22), Inches(9.5), Inches(0.62))
    tf = _tf(tb)
    para(tf, [(deck_name, {'size': 22, 'weight': 'eb'}),
              ("   |   ", {'size': 20, 'color': GRAY_TX}),
              (section_label, {'size': 20, 'weight': 'eb'})], first=True)
    hline(slide, MARGIN_X, 0.98, MARGIN_X + CONTENT_W, color=INK, width=1.5)


def add_headline(slide, text, y=ZONE_HEADLINE_Y):
    tb = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(y), Inches(CONTENT_W), Inches(0.6))
    tf = _tf(tb)
    para(tf, text, size=26, weight='eb', align=PP_ALIGN.CENTER, first=True)
    return tb


def add_card(slide, x, y, w, h, pill_text, pill_w=2.8):
    card = rect(slide, x, y, w, h, fill=WHITE, line_color=BORDER, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.035)
    pill = rect(slide, x + w / 2 - pill_w / 2, y - 0.22, pill_w, 0.44,
                fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = _tf(pill)
    para(tf, pill_text, size=14, weight='eb', color=WHITE, align=PP_ALIGN.CENTER, first=True)
    return card


def add_card_head(slide, x, y, w, text, sub=None, divider=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.75))
    tf = _tf(tb)
    para(tf, text, size=15, bold=True, align=PP_ALIGN.CENTER, first=True, line=1.15)
    if sub:
        para(tf, sub, size=15, bold=True, align=PP_ALIGN.CENTER, line=1.15)
    if divider:
        hline(slide, x + 0.35, y + 0.82, x + w - 0.35, color=BORDER, width=0.75)
    return tb


def add_chip(slide, x, y, w, h, text, fill=NAVY, size=12):
    sp = rect(slide, x, y, w, h, fill=fill)
    tf = _tf(sp)
    para(tf, text, size=size, weight='eb', color=WHITE, align=PP_ALIGN.CENTER,
         first=True, line=1.1)
    return sp


def add_chip_row(slide, x, y, w, chip_text, title, sub=None, chip_w=1.55, chip_h=0.52):
    """[NAVY 칩] + 우측 제목/보조설명 행"""
    add_chip(slide, x, y + 0.04, chip_w, chip_h, chip_text)
    tb = slide.shapes.add_textbox(Inches(x + chip_w + 0.18), Inches(y - 0.05),
                                  Inches(w - chip_w - 0.18), Inches(0.95))
    tf = _tf(tb, anchor=MSO_ANCHOR.TOP)
    para(tf, title, size=12.5, bold=True, first=True, line=1.15)
    if sub:
        para(tf, sub, size=10.5, color=GRAY_TX, line=1.15, space_before=2)
    return tb


def add_band(slide, x, y, w, text, h=0.45):
    sp = rect(slide, x, y, w, h, fill=SKY)
    tf = _tf(sp)
    para(tf, text, size=12.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, first=True)
    return sp


def add_takeaway(slide, text):
    sp = rect(slide, MARGIN_X, ZONE_TAKEAWAY_Y, CONTENT_W, 0.55,
              fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    tf = _tf(sp)
    para(tf, text, size=16, weight='eb', color=WHITE, align=PP_ALIGN.CENTER, first=True)
    return sp


def add_timeline(slide, x, y, w, points):
    """points: [(label, sublabel, done:bool, highlight:bool), ...]"""
    hline(slide, x, y, x + w, color=BORDER, width=1.5)
    n = len(points)
    step = w / (n - 1) if n > 1 else 0
    for i, (label, sub, done, hot) in enumerate(points):
        cx = x + step * i
        r = 0.09
        dot = rect(slide, cx - r, y - r, r * 2, r * 2,
                   fill=NAVY if done else WHITE,
                   line_color=None if done else NAVY, line_w=1.25,
                   shape=MSO_SHAPE.OVAL)
        if not done:
            dot.fill.solid()
            dot.fill.fore_color.rgb = SKY if not hot else WHITE
            dot.line.color.rgb = NAVY
        tb = slide.shapes.add_textbox(Inches(cx - 0.7), Inches(y + 0.13), Inches(1.4), Inches(0.55))
        tf = _tf(tb, anchor=MSO_ANCHOR.TOP)
        para(tf, label, size=11, bold=True,
             color=RED if hot else INK, align=PP_ALIGN.CENTER, first=True)
        if sub:
            para(tf, sub, size=9, color=GRAY_TX, align=PP_ALIGN.CENTER)


def add_entity(slide, x, y, w, h, title, sub=None, kind="plan"):
    """kind: legacy(STEEL 채움) / plan(NAVY 채움) / future(점선)"""
    if kind == "legacy":
        sp = rect(slide, x, y, w, h, fill=STEEL)
        tcol = WHITE
    elif kind == "plan":
        sp = rect(slide, x, y, w, h, fill=NAVY)
        tcol = WHITE
    else:
        sp = rect(slide, x, y, w, h, fill=WHITE, line_color=NAVY, line_w=1.0,
                  dash=MSO_LINE.DASH)
        tcol = NAVY
    tf = _tf(sp)
    para(tf, title, size=11.5, bold=True, color=tcol, align=PP_ALIGN.CENTER,
         first=True, line=1.1)
    if sub:
        para(tf, sub, size=9.5, color=tcol, align=PP_ALIGN.CENTER, line=1.1)
    return sp


def add_arrow(slide, x, y, w, h, direction="right", fill=NAVY):
    shape = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.UP_ARROW
    return rect(slide, x, y, w, h, fill=fill, shape=shape)


def add_track(slide, x, y, w, h, text, fill=NAVY, text_color=WHITE, dashed=False):
    """Two-Track 로드맵용 펜타곤 화살표"""
    if dashed:
        sp = rect(slide, x, y, w, h, fill=WHITE, line_color=NAVY, line_w=1.0,
                  dash=MSO_LINE.DASH)
        tc = NAVY
    else:
        sp = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.PENTAGON)
        tc = text_color
    tf = _tf(sp)
    para(tf, text, size=11.5, bold=True, color=tc, align=PP_ALIGN.CENTER, first=True)
    return sp


def caption(slide, x, y, w, text, align=PP_ALIGN.CENTER, color=GRAY_TX, size=10):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    tf = _tf(tb, anchor=MSO_ANCHOR.TOP)
    para(tf, text, size=size, color=color, align=align, first=True, line=1.15)
    return tb

# ──────────────────────────────── 4. 슬라이드 빌드 ────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


# ── S1. 커버 ─────────────────────────────────────────────
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
rect(s, 0, 0, 13.33, 0.18, fill=NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6))
tf = _tf(tb, anchor=MSO_ANCHOR.TOP)
para(tf, "Idea Pitch Template", size=44, weight='eb', color=NAVY, first=True)
para(tf, "Design System  |  토큰 · 컴포넌트 · 레이아웃 패턴", size=20, color=GRAY_TX,
     space_before=8)
# 컬러 토큰 스트립
swatches = [("NAVY", NAVY), ("BLUE", BLUE), ("SKY", SKY), ("STEEL", STEEL), ("RED", RED)]
for i, (name, c) in enumerate(swatches):
    x = 0.95 + i * 1.35
    rect(s, x, 4.35, 1.1, 0.55, fill=c, line_color=BORDER if name == "SKY" else None)
    caption(s, x, 4.95, 1.1, name)
caption(s, 0.95, 5.7, 11, "Header → Governing Headline → Section Cards → Takeaway Bar 의 4-Zone 구조",
        align=PP_ALIGN.LEFT, size=12, color=INK)
add_takeaway(s, "이 파일의 '빈 템플릿' 슬라이드를 복제하거나 generate_template.py 컴포넌트를 조합해 사용")

# ── S2. 디자인 토큰 ───────────────────────────────────────
s = new_slide()
add_header(s, "Design System", "① 디자인 토큰")
add_headline(s, "컬러는 NAVY 1색 지배 · 정보의 시제를 색·선 스타일로 구분")

card_l = add_card(s, MARGIN_X, ZONE_BODY_Y, 6.0, ZONE_BODY_H, "컬러 토큰")
tokens = [
    ("NAVY  #1F3864", "섹션 필 · 칩 · 결론 바 · 계획 개체", NAVY, WHITE),
    ("BLUE  #1273C4", "포인트 액션 · 링크", BLUE, WHITE),
    ("SKY  #DCE7F2", "강조 밴드 · 중간 결론", SKY, NAVY),
    ("STEEL  #8C96A0", "현황·레거시 개체", STEEL, WHITE),
    ("GRAY  #7F7F7F", "보조설명 · → 파생 문장", None, GRAY_TX),
    ("RED  #C00000", "'지금' 시점 마커 (슬라이드당 1개)", None, RED),
]
for i, (name, use, fill, tcol) in enumerate(tokens):
    yy = 2.30 + i * 0.68
    sp = rect(s, 0.85, yy, 2.1, 0.5, fill=fill if fill else WHITE,
              line_color=None if fill else BORDER)
    tf = _tf(sp)
    para(tf, name, size=10.5, bold=True, color=tcol, align=PP_ALIGN.CENTER, first=True)
    caption(s, 3.15, yy + 0.08, 3.1, use, align=PP_ALIGN.LEFT, color=INK, size=10.5)

card_r = add_card(s, 6.75, ZONE_BODY_Y, 6.13, ZONE_BODY_H, "타이포그래피")
types = [
    ("Headline  26pt 나눔스퀘어 ExtraBold", 20, 'eb', INK),
    ("Card Head  15pt 나눔스퀘어 Bold", 15, 'b', INK),
    ("Body  12pt 나눔스퀘어 Regular", 12, 'r', INK),
    ("Sub / 각주  10.5pt 나눔스퀘어 Light Gray", 10.5, 'l', GRAY_TX),
    ("Chip · Pill  12–14pt ExtraBold White (NAVY 배경)", 12, 'eb', NAVY),
    ("Takeaway  16pt ExtraBold White (NAVY 바)", 14, 'eb', NAVY),
]
tb = s.shapes.add_textbox(Inches(7.15), Inches(2.35), Inches(5.4), Inches(3.6))
tf = _tf(tb, anchor=MSO_ANCHOR.TOP)
first = True
for text, sz, wt, col in types:
    para(tf, text, size=sz, weight=wt, color=col, first=first, space_after=10, line=1.1)
    first = False
caption(s, 7.15, 6.05, 5.4, "기본 서체: 나눔스퀘어 NanumSquareOTF_ac — L/R/B/EB 4웨이트 (hangeul.naver.com)",
        align=PP_ALIGN.LEFT)
add_takeaway(s, "토큰 외 임의 색·크기 사용 금지 — 시제 규칙(과거 STEEL / 계획 NAVY / 미확정 점선)을 유지")

# ── S3. 레이아웃 아나토미 ─────────────────────────────────
s = new_slide()
add_header(s, "Design System", "② 레이아웃 구조")

zones = [
    (MARGIN_X, 0.18, CONTENT_W, 0.80, ""),
    (MARGIN_X, ZONE_HEADLINE_Y, CONTENT_W, 0.60, "② Governing Headline — 26pt Bold 중앙, 슬라이드의 '판단' 한 줄"),
    (MARGIN_X, ZONE_BODY_Y, 6.0, ZONE_BODY_H, "③-L Body 좌측 카드 (6.0in)"),
    (6.75, ZONE_BODY_Y, 6.13, ZONE_BODY_H, "③-R Body 우측 카드 스택 (6.13in)"),
    (MARGIN_X, ZONE_TAKEAWAY_Y, CONTENT_W, 0.55, ""),
]
for x, y, w, h, label in zones:
    rect(s, x, y, w, h, fill=None, line_color=BLUE, line_w=1.0, dash=MSO_LINE.DASH)
    if label:
        tb = s.shapes.add_textbox(Inches(x + 0.15), Inches(y + h / 2 - 0.18), Inches(w - 0.3), Inches(0.4))
        tf = _tf(tb)
        para(tf, label, size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER, first=True)
caption(s, 5.5, 0.48, 5.3, "① Header — 덱이름 | 섹션명 + 룰",
        align=PP_ALIGN.LEFT, color=BLUE, size=12)
caption(s, MARGIN_X, ZONE_TAKEAWAY_Y + 0.06, CONTENT_W,
        "④ Takeaway Bar — NAVY 바, '행동 함의' 결론 (Headline과 다른 문장)", color=BLUE, size=12)
caption(s, MARGIN_X, 6.62 - 0.55, 6.0, "좌우 마진 0.45in · 거터 0.30in", size=9.5)

# ── S4. 컴포넌트 라이브러리 ───────────────────────────────
s = new_slide()
add_header(s, "Design System", "③ 컴포넌트")
add_headline(s, "핵심 컴포넌트의 조합만으로 모든 슬라이드를 구성")

# 4-1 Section Card + Pill
add_card(s, 0.7, 2.35, 3.6, 1.45, "Section Pill Card", pill_w=2.5)
caption(s, 0.7, 3.9, 3.6, "add_card() — 필이 카드 상단 테두리에 걸침")
# 4-2 Chip row
add_chip_row(s, 4.85, 2.5, 4.35, "라벨 칩", "add_chip_row() — 제목",
             sub="→ 보조 설명은 GRAY 10.5pt")
caption(s, 4.85, 3.9, 3.6, "Label Chip + 설명 행")
# 4-4 Emphasis band
add_band(s, 0.7, 4.55, 3.6, "Emphasis Band — 카드 내 중간 결론")
caption(s, 0.7, 5.1, 3.6, "add_band()")
# 4-5 Timeline
add_timeline(s, 5.15, 4.8, 3.3, [
    ("완료", None, True, False),
    ("현재", None, False, True),
    ("예정", None, False, False),
])
caption(s, 4.85, 5.55, 3.9, "add_timeline() — 현재는 RED 라벨")
# 4-6 Entity boxes
add_entity(s, 9.2, 4.45, 1.35, 0.75, "현황", kind="legacy")
add_entity(s, 10.75, 4.45, 1.35, 0.75, "계획", kind="plan")
add_entity(s, 9.2, 5.35, 2.9, 0.55, "미확정 수요 (점선)", kind="future")
caption(s, 9.2, 5.95, 2.9, "add_entity() — 시제 3종")
# 4-7 Arrows
add_arrow(s, 0.9, 5.7, 1.0, 0.45, "right")
add_track(s, 2.1, 5.7, 2.1, 0.45, "Track Arrow")
caption(s, 0.9, 6.25, 3.3, "add_arrow() / add_track()")
add_takeaway(s, "Takeaway Bar — add_takeaway() : 슬라이드당 1개, 16pt Bold White")

# ── S5. 예시 패턴 A: 진단형 (상황 인식) ───────────────────
s = new_slide()
add_header(s, "Idea Pitch", "① 상황 인식")
add_headline(s, "핵심 판단 한 줄 — 현재 스탠스와 당사 구조의 특수성 제시")

# 좌측: 방향성 카드 (칩 행 3개)
add_card(s, MARGIN_X, ZONE_BODY_Y, 6.0, ZONE_BODY_H, "당사 방향성 (출처)")
add_card_head(s, MARGIN_X, 2.15, 6.0, "카드 서브헤드 — 현황 요약 한 줄")
rows = [
    ("주제 영역", "다양한 항목이 제안되고 있으나, 아직 초기 단계", "관련 사례 A, 사례 B, 사례 C 나열"),
    ("차세대 BM", "초기 아이디어 - 단, 현실적 제약 有", "제약 수치 1 · 제약 수치 2 · 미정 항목"),
    ("주요 Stance", "이해관계자별 상반된 입장 요약", "외부 권고 근거 ('34년~ 본격 투자)"),
]
for i, (chip_t, title, sub) in enumerate(rows):
    add_chip_row(s, 0.85, 3.25 + i * 0.95, 5.4, chip_t, title, sub=sub)
add_band(s, 0.85, 6.02, 5.2, "▶  중간 결론 — 방향성은 도출, 단 전제 조건은 미확정")

# 우측 상단: 타임라인 카드
add_card(s, 6.75, ZONE_BODY_Y, 6.13, 2.15, "생태계 동향 & Timeline")
add_card_head(s, 6.75, 2.12, 6.13, "기술 개발·표준화 진행 중 — 요구사항 반영의 창은 한시적",
              divider=False)
add_timeline(s, 7.35, 3.35, 4.9, [
    ("'23.6", "비전 수립", True, False),
    ("현재", "개발 경쟁", False, True),
    ("'29.3", "표준 배포", True, False),
    ("'31~", "상용화", False, False),
])
# 우측 하단: 구조 비교 카드
add_card(s, 6.75, 4.45, 6.13, 2.15, "당사 구조의 특수성")
add_card_head(s, 6.75, 4.62, 6.13, "글로벌 다수와 다른 '집중형 구조'", divider=False)
add_entity(s, 7.15, 5.45, 2.55, 0.85, "Global 다수", "분산형 구조", kind="legacy")
add_entity(s, 10.05, 5.45, 2.55, 0.85, "당사", "집중형 구조", kind="plan")
add_takeaway(s, "행동 함의 결론 — 당사 특수성을 고려한 기술 개발 및 생태계 조성이 중요함")

# ── S6. 예시 패턴 B: 제언형 (방향성 제언) ─────────────────
s = new_slide()
add_header(s, "Idea Pitch", "② 방향성 제언")
add_headline(s, "핵심 판단 한 줄 — 투자의 시계와 설계의 시계는 다름")

# 좌측: 확장성 카드 (전환 다이어그램)
add_card(s, MARGIN_X, ZONE_BODY_Y, 6.55, ZONE_BODY_H, "확장성 논거")
add_card_head(s, MARGIN_X, 2.15, 6.55, "당사 구조가 유리한 기반임을 보여주는 전환 다이어그램")
add_entity(s, 0.85, 3.25, 1.75, 1.75, "현재 자산", "(레거시)", kind="legacy")
add_arrow(s, 2.75, 3.9, 0.75, 0.45, "right")
caption(s, 2.6, 4.45, 1.05, "전환 기술", size=9)
add_entity(s, 3.65, 3.25, 2.0, 0.5, "+ 신규 수요분", kind="future")
add_entity(s, 3.65, 3.88, 2.0, 0.5, "+ 증가분", kind="future")
add_entity(s, 3.65, 4.5, 2.0, 0.5, "전환 개체 (계획)", kind="plan")
add_arrow(s, 5.95, 3.6, 0.5, 0.9, "up", fill=STEEL)
caption(s, 5.7, 4.55, 1.1, "수요 맞춤 증설", size=9)
add_band(s, 0.85, 5.25, 5.9, "자원 공유 기술과 결합 시, 비용 효율적·유연한 구축 가능")
add_chip(s, 0.85, 5.9, 1.3, 0.42, "하방 방어")
caption(s, 2.25, 5.95, 1.8, "수요 없어도 TCO 절감", align=PP_ALIGN.LEFT, color=INK, size=10)
add_chip(s, 4.0, 5.9, 1.3, 0.42, "상방 확장")
caption(s, 5.4, 5.95, 1.5, "수요 시 유연 증설", align=PP_ALIGN.LEFT, color=INK, size=10)

# 우측 상단: 한계 요인 카드
add_card(s, 7.3, ZONE_BODY_Y, 5.58, 2.35, "한계 요인")
add_chip_row(s, 7.6, 2.35, 5.0, "기울어진\n운동장", "외부 수요 중심 구조로 대기 시 선택지 소멸",
             sub="→ 수동적 Wait & see 시 선택지 소멸", chip_w=1.35, chip_h=0.62)
add_chip_row(s, 7.6, 3.3, 5.0, "절반의\n설계도", "Cross-domain 고려한 구조 준비 필요",
             sub="→ 단독 정의 불가 · 전사 공동 설계", chip_w=1.35, chip_h=0.62)
# 우측 하단: Two-Track 카드
add_card(s, 7.3, 4.65, 5.58, 1.95, "방향성 제언")
add_card_head(s, 7.3, 4.8, 5.58, "Two-Track — 투자는 신중히, 확보는 가속", divider=False)
caption(s, 8.75, 5.42, 1.2, "지금", color=RED, size=10)
add_chip(s, 7.6, 5.7, 1.15, 0.42, "투자", size=11)
add_track(s, 8.9, 5.7, 1.9, 0.42, "Wait & See", dashed=True)
add_track(s, 10.85, 5.7, 1.75, 0.42, "단계적 투자", fill=STEEL)
add_chip(s, 7.6, 6.2, 1.15, 0.42, "설계", size=11)
add_track(s, 8.9, 6.2, 3.7, 0.42, "전사 협력 — 확보 가속")
add_takeaway(s, "행동 함의 결론 — 투자는 신중히 하되, 규합은 전사 협력으로 적극 확장 필요")

# ── S7. 빈 템플릿 A: 1 + 2 구성 ──────────────────────────
s = new_slide()
add_header(s, "Idea Pitch", "① [섹션명]")
add_headline(s, "[거버닝 메시지 — 이 슬라이드의 판단 한 줄]")
add_card(s, MARGIN_X, ZONE_BODY_Y, 6.0, ZONE_BODY_H, "[좌측 섹션 라벨]")
add_card_head(s, MARGIN_X, 2.15, 6.0, "[카드 서브헤드]")
add_card(s, 6.75, ZONE_BODY_Y, 6.13, 2.15, "[우상단 섹션 라벨]")
add_card(s, 6.75, 4.45, 6.13, 2.15, "[우하단 섹션 라벨]")
add_takeaway(s, "[행동 함의 결론 — Headline과 다른 문장으로]")

# ── S8. 빈 템플릿 B: 2단 구성 ────────────────────────────
s = new_slide()
add_header(s, "Idea Pitch", "② [섹션명]")
add_headline(s, "[거버닝 메시지 — 이 슬라이드의 판단 한 줄]")
add_card(s, MARGIN_X, ZONE_BODY_Y, 6.05, ZONE_BODY_H, "[좌측 섹션 라벨]")
add_card(s, 6.83, ZONE_BODY_Y, 6.05, ZONE_BODY_H, "[우측 섹션 라벨]")
add_takeaway(s, "[행동 함의 결론]")

prs.save("idea_pitch_template.pptx")
print(f"저장 완료: idea_pitch_template.pptx ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
